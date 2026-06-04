#!/usr/bin/env python3
"""
pptx_tool.py - Inspect, check, and edit Microsoft PowerPoint .pptx files.

Subcommands:
  inspect <path>        - Full structural report of the deck.
  check <path>          - Detect common issues, including template text and empty placeholders.
  list-layouts <path>   - List all available slide layouts in the master.
  set-text <path> <slide_idx> <placeholder_name> <text>
                        - Set text in a named shape or placeholder on a slide.
  replace-text <path> <old_text> <new_text>
                        - Replace all occurrences of a text string globally.
  add-slide <path> <layout_name_substr>
                        - Add a new blank slide based on a matching layout.
  set-notes <path> <slide_idx> <text>
                        - Set speaker notes on a slide.
  export-text <path>    - Extract all slide text and notes as structured JSON.

Output: JSON to stdout. On error: JSON to stderr, exit 1.
"""

import sys, json, os, re, shutil, time
from pptx import Presentation

_TEMPLATE_PATTERNS = [
    re.compile(r'单击此处添加\s*(标题|副标题|文本|备注|页脚)', re.UNICODE),
    re.compile(r'Click\s+to\s+add\s+(title|subtitle|text|notes|footer)', re.IGNORECASE),
    re.compile(r'^\s*标题\s*$'),
    re.compile(r'^\s*副标题\s*$'),
    re.compile(r'^\s*正文\s*$'),
    re.compile(r'^\s*Text\s*$'),
    re.compile(r'^\s*Title\s*$'),
    re.compile(r'^\s*Subtitle\s*$'),
    re.compile(r'^\s*Body\s*$'),
    re.compile(r'^\s*Footer\s*$'),
    re.compile(r'添加备注.*$', re.UNICODE),
    re.compile(r'双击.*添加.*文本', re.UNICODE),
]


def _load(path):
    if not os.path.isfile(path):
        raise FileNotFoundError(f"File not found: {path}")
    return Presentation(path)


def _backup_before_write(path):
    backup = f"{path}.bak-{time.strftime('%Y%m%dT%H%M%S')}"
    shutil.copy2(path, backup)
    return os.path.abspath(backup)


def _slide_shape_info(shape):
    info = {"name": shape.name, "shape_type": str(shape.shape_type)}
    if shape.has_text_frame:
        para_texts = []
        for p in shape.text_frame.paragraphs:
            t = p.text.strip()
            if t:
                para_texts.append(t)
        info["text"] = para_texts
        info["text_joined"] = shape.text_frame.text.strip()
    if shape.has_table:
        rows = []
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        info["table"] = rows
    if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
        info["placeholder_idx"] = shape.placeholder_format.idx
        info["placeholder_type"] = str(shape.placeholder_format.type)
    return info


def _is_template_text(text):
    return any(p.search(text) for p in _TEMPLATE_PATTERNS)


def cmd_inspect(args):
    """inspect <path>"""
    if len(args) < 1:
        raise ValueError("Usage: inspect <path>")
    prs = _load(args[0])

    result = {
        "slide_width": str(prs.slide_width),
        "slide_height": str(prs.slide_height),
        "slide_count": len(prs.slides),
        "slides": [],
    }

    for idx, slide in enumerate(prs.slides, 1):
        sl = {"slide_index": idx, "layout_name": slide.slide_layout.name, "shapes": []}
        for shape in slide.shapes:
            sl["shapes"].append(_slide_shape_info(shape))
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                sl["notes"] = notes_text
        result["slides"].append(sl)

    result["layouts"] = [{"name": ly.name, "placeholders": [
        {"idx": ph.placeholder_format.idx, "name": ph.name, "type": str(ph.placeholder_format.type)}
        for ph in ly.placeholders
    ]} for ly in prs.slide_layouts]

    return {"ok": True, "data": result}


def cmd_check(args):
    """check <path>"""
    if len(args) < 1:
        raise ValueError("Usage: check <path>")
    prs = _load(args[0])

    issues = {"template_text_not_replaced": [], "empty_title_placeholder": [],
              "empty_body_placeholder": [], "notes_with_template_text": [],
              "slides_without_notes": []}

    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if _is_template_text(txt):
                    issues["template_text_not_replaced"].append({
                        "slide": idx, "shape": shape.name, "text": txt
                    })
                if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
                    ptype = str(shape.placeholder_format.type)
                    if ptype in ("TITLE", "CENTER_TITLE", "SUBTITLE") and not txt:
                        issues["empty_title_placeholder"].append({
                            "slide": idx, "shape": shape.name, "idx": shape.placeholder_format.idx
                        })
                    if ptype in ("BODY", "OBJECT", "TEXT") and not txt:
                        issues["empty_body_placeholder"].append({
                            "slide": idx, "shape": shape.name, "idx": shape.placeholder_format.idx
                        })

        if slide.has_notes_slide:
            notes_txt = slide.notes_slide.notes_text_frame.text.strip()
            if _is_template_text(notes_txt):
                issues["notes_with_template_text"].append({
                    "slide": idx, "text": notes_txt
                })
        else:
            issues["slides_without_notes"].append({"slide": idx})

    result = {k: v for k, v in issues.items() if v}

    issue_count = sum(len(v) for v in result.values())
    if issue_count:
        summary = f"Checked {len(prs.slides)} slides and found {issue_count} issue(s)."
    else:
        summary = f"Checked {len(prs.slides)} slides and found no issues."

    return {"ok": True, "data": {"summary": summary, "issues": result}}


def cmd_list_layouts(args):
    """list-layouts <path>"""
    if len(args) < 1:
        raise ValueError("Usage: list-layouts <path>")
    prs = _load(args[0])
    layouts = []
    for ly in prs.slide_layouts:
        layouts.append({
            "name": ly.name,
            "placeholders": [
                {"idx": ph.placeholder_format.idx, "name": ph.name,
                 "type": str(ph.placeholder_format.type)}
                for ph in ly.placeholders
            ]
        })
    return {"ok": True, "data": {"layouts": layouts, "count": len(layouts)}}


def cmd_set_text(args):
    """set-text <path> <slide_idx> <placeholder_name_substr> <text>"""
    if len(args) < 4:
        raise ValueError("Usage: set-text <path> <slide_idx> <placeholder_name_substr> <text>")
    path, slide_idx_str, ph_name_sub, new_text = args[0], args[1], args[2], args[3]
    slide_idx = int(slide_idx_str) - 1
    prs = _load(path)
    if slide_idx < 0 or slide_idx >= len(prs.slides):
        raise ValueError(f"Slide index out of range (1-{len(prs.slides)})")

    slide = prs.slides[slide_idx]
    found = False
    for shape in slide.shapes:
        if ph_name_sub.lower() in shape.name.lower() and shape.has_text_frame:
            tf = shape.text_frame
            for para in tf.paragraphs:
                para.clear()
            tf.paragraphs[0].text = new_text
            found = True
            break

    if not found:
        raise ValueError(f"No text-enabled shape matching '{ph_name_sub}' found on slide {slide_idx_str}")

    backup_path = _backup_before_write(path)
    prs.save(path)
    return {"ok": True, "data": {"action": f"Set text on slide {slide_idx_str} shape matching '{ph_name_sub}'", "path": os.path.abspath(path), "backup_path": backup_path}}


def cmd_replace_text(args):
    """replace-text <path> <old_text> <new_text>"""
    if len(args) < 3:
        raise ValueError("Usage: replace-text <path> <old_text> <new_text>")
    path, old_text, new_text = args[0], args[1], args[2]
    prs = _load(path)
    replacements = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                new_shape_text = tf.text.replace(old_text, new_text)
                if new_shape_text != tf.text:
                    for para in tf.paragraphs:
                        para.clear()
                    tf.paragraphs[0].text = new_shape_text
                    replacements += 1
    backup_path = _backup_before_write(path)
    prs.save(path)
    return {"ok": True, "data": {"replacements": replacements, "old_text": old_text, "new_text": new_text, "path": os.path.abspath(path), "backup_path": backup_path}}


def cmd_add_slide(args):
    """add-slide <path> <layout_name_substr>"""
    if len(args) < 2:
        raise ValueError("Usage: add-slide <path> <layout_name_substr>")
    path, ly_name_sub = args[0], args[1]
    prs = _load(path)

    found_layout = None
    for layout in prs.slide_layouts:
        if ly_name_sub.lower() in layout.name.lower():
            found_layout = layout
            break
    if not found_layout:
        available = ", ".join([f"'{l.name}'" for l in prs.slide_layouts[:5]])
        if len(prs.slide_layouts) > 5:
            available += ", ..."
        raise ValueError(f"No layout matching '{ly_name_sub}' found. Available: {available}")

    new_idx = len(prs.slides) + 1
    prs.slides.add_slide(found_layout)
    backup_path = _backup_before_write(path)
    prs.save(path)
    return {"ok": True, "data": {"action": f"Added new slide {new_idx} using layout '{found_layout.name}'", "new_slide_index": new_idx, "layout": found_layout.name, "path": os.path.abspath(path), "backup_path": backup_path}}


def cmd_set_notes(args):
    """set-notes <path> <slide_idx> <text>"""
    if len(args) < 3:
        raise ValueError("Usage: set-notes <path> <slide_idx> <text>")
    path, slide_idx_str, new_text = args[0], args[1], args[2]
    slide_idx = int(slide_idx_str) - 1
    prs = _load(path)
    if slide_idx < 0 or slide_idx >= len(prs.slides):
        raise ValueError(f"Slide index out of range (1-{len(prs.slides)})")

    slide = prs.slides[slide_idx]
    notes_slide = slide.notes_slide
    if notes_slide is None:
        try:
            notes_slide = slide.notes_slide
        except Exception:
            pass
    if notes_slide and notes_slide.notes_text_frame:
        tf = notes_slide.notes_text_frame
        for para in tf.paragraphs:
            para.clear()
        tf.paragraphs[0].text = new_text
    else:
        raise ValueError(f"Cannot add notes to slide {slide_idx_str}; notes slide does not exist")

    backup_path = _backup_before_write(path)
    prs.save(path)
    return {"ok": True, "data": {"action": f"Set notes on slide {slide_idx_str}", "path": os.path.abspath(path), "backup_path": backup_path}}


def cmd_export_text(args):
    """export-text <path>"""
    if len(args) < 1:
        raise ValueError("Usage: export-text <path>")
    prs = _load(args[0])
    slides_out = []
    for idx, slide in enumerate(prs.slides, 1):
        entry = {"slide_index": idx, "layout": slide.slide_layout.name}
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    texts.append({"shape": shape.name, "text": txt})
        entry["texts"] = texts
        if slide.has_notes_slide:
            nt = slide.notes_slide.notes_text_frame.text.strip()
            if nt:
                entry["notes"] = nt
        slides_out.append(entry)

    return {"ok": True, "data": {"slide_count": len(slides_out), "slides": slides_out}}


COMMANDS = {
    "inspect": cmd_inspect,
    "check": cmd_check,
    "list-layouts": cmd_list_layouts,
    "set-text": cmd_set_text,
    "replace-text": cmd_replace_text,
    "add-slide": cmd_add_slide,
    "set-notes": cmd_set_notes,
    "export-text": cmd_export_text,
}


def main(argv):
    if len(argv) < 1 or argv[0] in ("-h", "--help", "help"):
        print("Usage: pptx_tool.py <command> [args...]\n")
        print("Commands:")
        for name, fn in COMMANDS.items():
            print(f"  {name:<15} {fn.__doc__.strip().split(chr(10))[0]}")
        return {"ok": True}

    cmd = argv[0]
    if cmd not in COMMANDS:
        raise ValueError(f"Unknown command: {cmd}. Available: {', '.join(COMMANDS)}")

    return COMMANDS[cmd](argv[1:])


if __name__ == "__main__":
    try:
        result = main(sys.argv[1:])
    except Exception as e:
        print(json.dumps({"ok": False, "error": str(e)}, ensure_ascii=False), file=sys.stderr)
        sys.exit(1)
    print(json.dumps(result, ensure_ascii=False))
